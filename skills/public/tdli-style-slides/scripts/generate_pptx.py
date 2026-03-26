#!/usr/bin/env python3
"""Generate TDLI-branded PPTX from a JSON content plan.

Fallback generator for environments without xelatex.
Produces a PPTX file with TDLI branding (cover, header/footer, colors, logo)
and supports 3 layout types matching the LaTeX Beamer template.

Usage:
    python generate_pptx.py --plan-file plan.json --output-file output.pptx \
        --assets-dir /path/to/assets

The assets-dir should contain: 封面.png, 页眉newstyle.png, 页脚.png,
校标-标志中英文横版.png
"""

import argparse
import json
import os
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------------------
# Brand constants
# ---------------------------------------------------------------------------
SJTU_BLUE = RGBColor(0x11, 0x4A, 0x79)
SJTU_LIGHT_BLUE = RGBColor(0x00, 0x71, 0xBC)
ACCENT_GOLD = RGBColor(0xB4, 0x8C, 0x32)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x33, 0x33, 0x33)
GRAY_TEXT = RGBColor(0x99, 0x99, 0x99)
LIGHT_BG = RGBColor(0xF5, 0xF8, 0xFC)
BLOCK_TITLE_BG = SJTU_BLUE
TAG_BG = RGBColor(0xD6, 0xE4, 0xF0)  # sjtuBlue 15%

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Layout geometry
HEADER_HEIGHT = Inches(0.85)
FOOTER_HEIGHT = Inches(0.55)
CONTENT_TOP = Inches(1.05)
CONTENT_BOTTOM = SLIDE_HEIGHT - FOOTER_HEIGHT - Inches(0.15)
CONTENT_LEFT = Inches(0.6)
CONTENT_RIGHT = SLIDE_WIDTH - Inches(0.6)
CONTENT_WIDTH = CONTENT_RIGHT - CONTENT_LEFT
CONTENT_HEIGHT = CONTENT_BOTTOM - CONTENT_TOP

FONT_TITLE = "Arial"
FONT_BODY = "Arial"
FONT_CJK = "Microsoft YaHei"  # fallback CJK font available on most systems


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def _add_image_safe(slide, img_path, left, top, width, height=None):
    """Add image to slide if file exists, skip silently otherwise."""
    if os.path.exists(img_path):
        if height:
            slide.shapes.add_picture(img_path, left, top, width, height)
        else:
            slide.shapes.add_picture(img_path, left, top, width)
        return True
    print(f"Warning: image not found, skipping: {img_path}", file=sys.stderr)
    return False


def _add_textbox(slide, left, top, width, height):
    """Add a textbox and return its text_frame."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf


def _set_paragraph(
    tf,
    text,
    font_size=12,
    bold=False,
    color=DARK_TEXT,
    alignment=PP_ALIGN.LEFT,
    font_name=FONT_BODY,
    space_after=Pt(4),
):
    """Set text on the first paragraph of a text_frame."""
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = space_after
    return p


def _add_paragraph(
    tf,
    text,
    font_size=12,
    bold=False,
    color=DARK_TEXT,
    alignment=PP_ALIGN.LEFT,
    font_name=FONT_BODY,
    space_after=Pt(4),
    level=0,
):
    """Append a new paragraph to a text_frame."""
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = space_after
    p.level = level
    return p


def _add_block(slide, left, top, width, title, items, tags=None, metrics=None):
    """Add a styled block (title bar + content) to a slide. Returns bottom Y."""
    title_height = Inches(0.35)
    item_line_height = Inches(0.28)
    tag_line_height = Inches(0.25) if tags else 0
    padding = Inches(0.1)

    n_items = len(items) if items else 0
    n_metrics = len(metrics) if metrics else 0
    body_height = (
        (n_items + n_metrics) * item_line_height + tag_line_height + padding * 2
    )
    total_height = title_height + body_height

    # Title bar
    title_shape = slide.shapes.add_shape(
        1,
        left,
        top,
        width,
        title_height,  # MSO_SHAPE.RECTANGLE = 1
    )
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = BLOCK_TITLE_BG
    title_shape.line.fill.background()
    tf = title_shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    _set_paragraph(tf, title, font_size=11, bold=True, color=WHITE)

    # Body background
    body_top = top + title_height
    body_shape = slide.shapes.add_shape(1, left, body_top, width, body_height)
    body_shape.fill.solid()
    body_shape.fill.fore_color.rgb = LIGHT_BG
    body_shape.line.fill.background()

    # Body text
    tf = _add_textbox(
        slide,
        left + Inches(0.15),
        body_top + padding,
        width - Inches(0.3),
        body_height - padding * 2,
    )

    first = True

    # Tags line
    if tags:
        tag_text = "  ".join(f"[{t}]" for t in tags)
        if first:
            _set_paragraph(
                tf, tag_text, font_size=9, color=SJTU_BLUE, space_after=Pt(6)
            )
            first = False
        else:
            _add_paragraph(
                tf, tag_text, font_size=9, color=SJTU_BLUE, space_after=Pt(6)
            )

    # Items
    if items:
        for item in items:
            text = f"• {item}"
            if first:
                _set_paragraph(
                    tf, text, font_size=10, color=DARK_TEXT, space_after=Pt(3)
                )
                first = False
            else:
                _add_paragraph(
                    tf, text, font_size=10, color=DARK_TEXT, space_after=Pt(3)
                )

    # Metrics
    if metrics:
        for m in metrics:
            text = f"▸ {m}"
            if first:
                _set_paragraph(
                    tf,
                    text,
                    font_size=10,
                    bold=True,
                    color=SJTU_BLUE,
                    space_after=Pt(3),
                )
                first = False
            else:
                _add_paragraph(
                    tf,
                    text,
                    font_size=10,
                    bold=True,
                    color=SJTU_BLUE,
                    space_after=Pt(3),
                )

    return top + total_height


def _add_header_footer(slide, assets_dir):
    """Add header and footer overlay images to a content slide."""
    header_img = os.path.join(assets_dir, "页眉newstyle.png")
    footer_img = os.path.join(assets_dir, "页脚.png")
    _add_image_safe(slide, header_img, 0, 0, SLIDE_WIDTH, HEADER_HEIGHT)
    _add_image_safe(
        slide, footer_img, 0, SLIDE_HEIGHT - FOOTER_HEIGHT, SLIDE_WIDTH, FOOTER_HEIGHT
    )


def _add_frame_title(slide, title):
    """Add white frame title text overlaid on the header image."""
    tf = _add_textbox(
        slide, Inches(0.5), Inches(0.15), SLIDE_WIDTH - Inches(1.0), Inches(0.55)
    )
    _set_paragraph(
        tf, title, font_size=18, bold=True, color=WHITE, font_name=FONT_TITLE
    )


def _add_speaker_notes(slide, notes_text):
    """Add speaker notes to a slide."""
    if not notes_text:
        return
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    if tf is not None:
        tf.text = notes_text


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------


def build_cover_slide(prs, metadata, assets_dir):
    """Build the TDLI cover page."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Background
    cover_img = os.path.join(assets_dir, "封面.png")
    _add_image_safe(slide, cover_img, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)

    # Logo
    logo_img = os.path.join(assets_dir, "校标-标志中英文横版.png")
    logo_width = Inches(3.2)
    logo_left = (SLIDE_WIDTH - logo_width) // 2
    _add_image_safe(slide, logo_img, logo_left, Inches(1.0), logo_width)

    # Photo (optional — only if user provides a photo path in metadata)
    photo_path = metadata.get("photo")
    if photo_path and os.path.exists(photo_path):
        slide.shapes.add_picture(
            photo_path, Inches(1.2), Inches(1.8), Inches(2.0), Inches(2.0)
        )

    # Name
    center_left = (
        Inches(3.5) if photo_path and os.path.exists(photo_path) else Inches(2.5)
    )
    center_width = Inches(6.5)

    tf = _add_textbox(slide, center_left, Inches(1.8), center_width, Inches(0.5))
    _set_paragraph(
        tf,
        metadata.get("presenter_name", ""),
        font_size=24,
        bold=True,
        color=SJTU_BLUE,
        alignment=PP_ALIGN.CENTER,
    )

    # Subtitle
    tf = _add_textbox(slide, center_left, Inches(2.35), center_width, Inches(0.35))
    _set_paragraph(
        tf,
        metadata.get("subtitle", ""),
        font_size=12,
        color=GRAY_TEXT,
        alignment=PP_ALIGN.CENTER,
    )

    # Cover summary box
    cover_lines = metadata.get("cover_lines", [])
    if cover_lines:
        box_top = Inches(3.0)
        box_width = Inches(7.0)
        box_left = (SLIDE_WIDTH - box_width) // 2
        box_height = Inches(0.3) * len(cover_lines) + Inches(0.4)

        # Border rectangle
        rect = slide.shapes.add_shape(1, box_left, box_top, box_width, box_height)
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        rect.line.color.rgb = SJTU_BLUE
        rect.line.width = Pt(1.5)

        tf = _add_textbox(
            slide,
            box_left + Inches(0.2),
            box_top + Inches(0.15),
            box_width - Inches(0.4),
            box_height - Inches(0.3),
        )
        for i, line in enumerate(cover_lines):
            if i == 0:
                _set_paragraph(
                    tf,
                    line,
                    font_size=10,
                    color=SJTU_BLUE,
                    alignment=PP_ALIGN.CENTER,
                    space_after=Pt(4),
                )
            else:
                _add_paragraph(
                    tf,
                    line,
                    font_size=10,
                    color=SJTU_BLUE,
                    alignment=PP_ALIGN.CENTER,
                    space_after=Pt(4),
                )

    # Affiliation + Date
    info_top = Inches(4.8) if cover_lines else Inches(3.5)
    tf = _add_textbox(slide, Inches(2.0), info_top, Inches(9.0), Inches(0.6))
    _set_paragraph(
        tf,
        metadata.get("affiliation", ""),
        font_size=10,
        color=DARK_TEXT,
        alignment=PP_ALIGN.CENTER,
    )
    _add_paragraph(
        tf,
        metadata.get("date", ""),
        font_size=10,
        color=DARK_TEXT,
        alignment=PP_ALIGN.CENTER,
    )

    # Contact info at bottom
    contact_parts = []
    if metadata.get("email"):
        contact_parts.append(f"✉ {metadata['email']}")
    if metadata.get("phone"):
        contact_parts.append(f"☎ {metadata['phone']}")
    if metadata.get("github"):
        contact_parts.append(f"⌨ {metadata['github']}")
    if contact_parts:
        tf = _add_textbox(slide, Inches(1.0), Inches(6.8), Inches(11.0), Inches(0.35))
        _set_paragraph(
            tf,
            "    ".join(contact_parts),
            font_size=8,
            color=GRAY_TEXT,
            alignment=PP_ALIGN.CENTER,
        )

    _add_speaker_notes(slide, metadata.get("cover_notes", ""))


def build_content_slide_a(prs, slide_data, assets_dir):
    """Layout A: Left text blocks + Right figure."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_footer(slide, assets_dir)
    _add_frame_title(slide, slide_data.get("title", ""))

    # Left column: blocks
    left_x = CONTENT_LEFT
    col_width = (CONTENT_WIDTH - Inches(0.3)) * 0.48
    current_y = CONTENT_TOP

    for block in slide_data.get("blocks", []):
        current_y = _add_block(
            slide,
            left_x,
            current_y,
            col_width,
            title=block.get("title", ""),
            items=block.get("items", []),
            tags=block.get("tags"),
            metrics=block.get("metrics"),
        )
        current_y += Inches(0.15)

    # Right column: figure
    right_x = CONTENT_LEFT + col_width + Inches(0.3)
    fig_width = (CONTENT_WIDTH - Inches(0.3)) * 0.49
    figure_path = slide_data.get("figure", "")
    if figure_path and os.path.exists(figure_path):
        _add_image_safe(
            slide,
            figure_path,
            right_x,
            CONTENT_TOP,
            fig_width,
            CONTENT_HEIGHT - Inches(0.4),
        )

    caption = slide_data.get("figure_caption", "")
    if caption:
        tf = _add_textbox(
            slide, right_x, CONTENT_BOTTOM - Inches(0.3), fig_width, Inches(0.25)
        )
        _set_paragraph(
            tf, caption, font_size=8, color=GRAY_TEXT, alignment=PP_ALIGN.CENTER
        )

    _add_speaker_notes(slide, slide_data.get("notes", ""))


def build_content_slide_b(prs, slide_data, assets_dir):
    """Layout B: Full-width stacked blocks."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_footer(slide, assets_dir)
    _add_frame_title(slide, slide_data.get("title", ""))

    current_y = CONTENT_TOP
    full_width = CONTENT_WIDTH

    for block in slide_data.get("blocks", []):
        current_y = _add_block(
            slide,
            CONTENT_LEFT,
            current_y,
            full_width,
            title=block.get("title", ""),
            items=block.get("items", []),
            tags=block.get("tags"),
            metrics=block.get("metrics"),
        )
        current_y += Inches(0.15)

    # Optional table
    table_data = slide_data.get("table")
    if table_data:
        headers = table_data.get("headers", [])
        rows = table_data.get("rows", [])
        if headers and rows:
            n_rows = len(rows) + 1  # +1 for header
            n_cols = len(headers)
            row_height = Inches(0.3)
            table_height = row_height * n_rows
            table_width = min(full_width, Inches(8.0))
            table_left = CONTENT_LEFT + (full_width - table_width) // 2

            table_shape = slide.shapes.add_table(
                n_rows,
                n_cols,
                table_left,
                current_y + Inches(0.1),
                table_width,
                table_height,
            )
            table = table_shape.table

            # Header row
            for j, h in enumerate(headers):
                cell = table.cell(0, j)
                cell.text = h
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(9)
                    p.font.bold = True
                    p.font.color.rgb = WHITE
                    p.alignment = PP_ALIGN.CENTER
                cell.fill.solid()
                cell.fill.fore_color.rgb = SJTU_BLUE

            # Data rows
            for i, row in enumerate(rows):
                for j, val in enumerate(row):
                    cell = table.cell(i + 1, j)
                    cell.text = str(val)
                    for p in cell.text_frame.paragraphs:
                        p.font.size = Pt(9)
                        p.font.color.rgb = DARK_TEXT
                        p.alignment = PP_ALIGN.CENTER

    _add_speaker_notes(slide, slide_data.get("notes", ""))


def build_content_slide_c(prs, slide_data, assets_dir):
    """Layout C: Two equal-width block columns."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_footer(slide, assets_dir)
    _add_frame_title(slide, slide_data.get("title", ""))

    col_width = (CONTENT_WIDTH - Inches(0.3)) * 0.49
    gap = Inches(0.3)

    # Left block
    left_block = slide_data.get("left_block", {})
    if left_block:
        _add_block(
            slide,
            CONTENT_LEFT,
            CONTENT_TOP,
            col_width,
            title=left_block.get("title", ""),
            items=left_block.get("items", []),
            tags=left_block.get("tags"),
            metrics=left_block.get("metrics"),
        )

    # Right block
    right_block = slide_data.get("right_block", {})
    if right_block:
        _add_block(
            slide,
            CONTENT_LEFT + col_width + gap,
            CONTENT_TOP,
            col_width,
            title=right_block.get("title", ""),
            items=right_block.get("items", []),
            tags=right_block.get("tags"),
            metrics=right_block.get("metrics"),
        )

    _add_speaker_notes(slide, slide_data.get("notes", ""))


def build_thankyou_slide(prs, metadata, assets_dir):
    """Build the thank-you page."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    cover_img = os.path.join(assets_dir, "封面.png")
    _add_image_safe(slide, cover_img, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)

    # Thank You text
    tf = _add_textbox(slide, Inches(2.0), Inches(2.5), Inches(9.0), Inches(1.0))
    _set_paragraph(
        tf,
        "Thank You",
        font_size=44,
        bold=True,
        color=SJTU_BLUE,
        alignment=PP_ALIGN.CENTER,
    )

    tf = _add_textbox(slide, Inches(2.0), Inches(3.5), Inches(9.0), Inches(0.5))
    _set_paragraph(
        tf,
        "Questions and Discussion",
        font_size=20,
        color=GRAY_TEXT,
        alignment=PP_ALIGN.CENTER,
    )

    # Contact info
    contact_parts = []
    if metadata.get("email"):
        contact_parts.append(f"✉  {metadata['email']}")
    if metadata.get("phone"):
        contact_parts.append(f"☎  {metadata['phone']}")
    if metadata.get("github"):
        contact_parts.append(f"⌨  {metadata['github']}")
    if contact_parts:
        tf = _add_textbox(slide, Inches(3.0), Inches(4.5), Inches(7.0), Inches(1.2))
        for i, part in enumerate(contact_parts):
            if i == 0:
                _set_paragraph(
                    tf,
                    part,
                    font_size=12,
                    color=DARK_TEXT,
                    alignment=PP_ALIGN.CENTER,
                    space_after=Pt(6),
                )
            else:
                _add_paragraph(
                    tf,
                    part,
                    font_size=12,
                    color=DARK_TEXT,
                    alignment=PP_ALIGN.CENTER,
                    space_after=Pt(6),
                )

    _add_speaker_notes(slide, metadata.get("thankyou_notes", ""))


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

LAYOUT_BUILDERS = {
    "A": build_content_slide_a,
    "B": build_content_slide_b,
    "C": build_content_slide_c,
}


def generate_pptx(plan_file: str, output_file: str, assets_dir: str) -> str:
    """Generate a TDLI-branded PPTX from a JSON content plan.

    Args:
        plan_file: Path to JSON content plan.
        output_file: Path to output PPTX file.
        assets_dir: Path to directory containing TDLI brand assets.

    Returns:
        Status message.
    """
    with open(plan_file, "r", encoding="utf-8") as f:
        plan = json.load(f)

    metadata = plan.get("metadata", {})
    slides = plan.get("slides", [])

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Cover page
    build_cover_slide(prs, metadata, assets_dir)

    # Content slides
    for slide_data in slides:
        layout = slide_data.get("layout", "B").upper()
        builder = LAYOUT_BUILDERS.get(layout)
        if builder:
            builder(prs, slide_data, assets_dir)
        else:
            print(
                f"Warning: unknown layout '{layout}', falling back to B",
                file=sys.stderr,
            )
            build_content_slide_b(prs, slide_data, assets_dir)

    # Thank you page
    build_thankyou_slide(prs, metadata, assets_dir)

    # Save
    os.makedirs(os.path.dirname(os.path.abspath(output_file)), exist_ok=True)
    prs.save(output_file)

    total_slides = len(slides) + 2  # +cover +thankyou
    return f"Generated TDLI presentation: {total_slides} slides → {output_file}"


if __name__ == "__main__":
    parser = argparse.ArgumentParser(
        description="Generate TDLI-branded PPTX from JSON content plan"
    )
    parser.add_argument("--plan-file", required=True, help="Path to JSON content plan")
    parser.add_argument("--output-file", required=True, help="Output PPTX file path")
    parser.add_argument(
        "--assets-dir", required=True, help="Path to TDLI brand assets directory"
    )

    args = parser.parse_args()

    try:
        print(generate_pptx(args.plan_file, args.output_file, args.assets_dir))
    except Exception as e:
        print(f"Error: {e}", file=sys.stderr)
        sys.exit(1)
